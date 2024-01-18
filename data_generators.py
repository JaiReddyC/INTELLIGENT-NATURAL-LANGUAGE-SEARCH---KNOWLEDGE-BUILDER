 data_generators.py
import streamlit as st
import spacy
from transformers import BartForConditionalGeneration, BartTokenizer
import pptx
import os
import pandas as pd
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from docx import Document
import PyPDF2
# Load the English language model
nlp = spacy.load("en_core_web_sm")
# Function to extract text from a ppt
def extract_text_from_ppt(file_path):
    try:
        presentation = Presentation(file_path)
        text = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text += paragraph.text + ' '
        return text
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return ''

# Function to extract text from different sources based on provided links DataFrame
# def generate_data_from_files(file_paths):
#     text_from_ppt = ""
#     for ppt_path in file_paths:
#         text_from_ppt += extract_text_from_ppt(ppt_path)
#     #nlp = spacy.load("en_core_web_sm")
#     doc = nlp(text_from_ppt)
#     knowledge_graph = {}
#     for entity in doc.ents:
#         if entity.label_:
#             knowledge_graph[entity.text] = {'label': entity.label_, 'text': entity.sent.text}
#     return knowledge_graph



# Function to extract text from a webpage or load from a single cached file
def extract_text_from_webpage(url):
    cache_file = "webpage_cache.txt"
    try:
        response = requests.get(url)
        soup = BeautifulSoup(response.content, 'html.parser')
        text = ' '.join([p.get_text() for p in soup.find_all('p')])
        # Cache the extracted text by appending it to the single file
        with open(cache_file, 'a', encoding='utf-8') as file:
            file.write(text + '\n\n')
        return text
    except Exception as e:
        print(f"Error extracting text from {url}: {e}")
        return ''


def extract_text_from_ppt(uploaded_file):
    presentation = Presentation(uploaded_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + " "
    return text

def extract_text_from_pdf(uploaded_file):
    pdf_reader = PyPDF2.PdfReader(uploaded_file)
    num_pages = len(pdf_reader.pages)
    text = ""
    for page_num in range(num_pages):
        page = pdf_reader.pages[page_num]
        text += page.extract_text()
    return text

def extract_text_from_docx(uploaded_file):
    text = ""
    try:
        # Load the uploaded document
        doc = Document(uploaded_file)
        # Extract text from each paragraph in the document
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
    return text
